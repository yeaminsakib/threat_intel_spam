from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.util import Inches, Pt


OUTPUT_PATH = Path("presentation/AI_Threat_Intelligence_Proposal.pptx")
TITLE_FONT = "Calibri"
BODY_FONT = "Calibri"


def style_title(shape):
    p = shape.text_frame.paragraphs[0]
    if p.runs:
        run = p.runs[0]
    else:
        run = p.add_run()
    run.font.name = TITLE_FONT
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(20, 40, 80)


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "AI-driven Threat Intelligence System for Early Detection of Phishing Campaigns"
    style_title(slide.shapes.title)

    subtitle = slide.placeholders[1]
    subtitle.text = (
        "Using Spam Email Analysis, IOC Correlation, and Machine Learning\n\n"
        "Student: ____________________\n"
        "Department / University: ____________________\n"
        "Supervisor: ____________________"
    )
    tf = subtitle.text_frame
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.name = BODY_FONT
            run.font.size = Pt(18)
            run.font.color.rgb = RGBColor(55, 55, 55)


def add_bullet_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    style_title(slide.shapes.title)

    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.clear()

    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0

        paragraph = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        paragraph.text = text
        paragraph.level = level
        paragraph.space_after = Pt(6)

        for run in paragraph.runs:
            run.font.name = BODY_FONT
            run.font.size = Pt(20 if level == 0 else 18)
            run.font.color.rgb = RGBColor(45, 45, 45)


def add_architecture_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "11. System Architecture"
    style_title(slide.shapes.title)

    steps = [
        "Spam Emails",
        "Processing",
        "IOC Extraction",
        "ML Model",
        "Threat Graph",
        "Prediction Output",
    ]

    left = Inches(0.4)
    top = Inches(2.1)
    box_w = Inches(1.5)
    box_h = Inches(1.0)
    gap = Inches(0.12)

    boxes = []
    for i, step in enumerate(steps):
        x = left + i * (box_w + gap)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, box_w, box_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(232, 242, 255)
        shape.line.color.rgb = RGBColor(80, 120, 180)

        text_frame = shape.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = step
        p.alignment = 1
        for run in p.runs:
            run.font.name = BODY_FONT
            run.font.size = Pt(15)
            run.font.bold = True
            run.font.color.rgb = RGBColor(30, 60, 100)
        boxes.append(shape)

    for i in range(len(boxes) - 1):
        s = boxes[i]
        e = boxes[i + 1]
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            s.left + s.width,
            s.top + s.height // 2,
            e.left,
            e.top + e.height // 2,
        )
        connector.line.color.rgb = RGBColor(80, 120, 180)
        connector.line.width = Pt(1.8)

    note_box = slide.shapes.add_textbox(Inches(0.7), Inches(4.0), Inches(8.8), Inches(1.6))
    tf = note_box.text_frame
    tf.clear()

    lines = [
        "• Pipeline supports batch analysis now and real-time SOC integration later.",
        "• Modular blocks allow independent updates for extraction, ML, and graph analytics.",
    ]
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        for run in p.runs:
            run.font.name = BODY_FONT
            run.font.size = Pt(18)
            run.font.color.rgb = RGBColor(45, 45, 45)


def build_presentation():
    prs = Presentation()

    add_title_slide(prs)

    slides = [
        (
            "2. Introduction",
            [
                "Spam email remains a major entry point for phishing attacks.",
                "Traditional filters focus on blocking messages, not understanding campaigns.",
                "Security teams need early warning from spam trends and attacker behavior.",
                "This project shifts from spam blocking to threat intelligence generation.",
            ],
        ),
        (
            "3. Problem Statement",
            [
                "Current systems mostly detect known spam patterns.",
                "Unknown and fast-changing phishing campaigns can bypass static rules.",
                "Spam data is rarely transformed into correlated intelligence.",
                "Result: delayed response and higher operational risk for SOC teams.",
            ],
        ),
        (
            "4. Motivation of the Project",
            [
                "Cyber attacks and phishing campaigns are increasing every year.",
                "Attackers quickly rotate domains, URLs, and sender identities.",
                "Organizations need predictive, not only reactive, defense.",
                "SOC teams require campaign-level context for faster response.",
            ],
        ),
        (
            "5. Related Work",
            [
                "Major platforms (e.g., Google, Microsoft) provide strong spam filtering.",
                "SIEM and TIP tools aggregate alerts but are often not spam-centered.",
                "Research explores phishing classification using NLP and ML.",
                "Gap in prior work: limited campaign prediction and weak IOC correlation.",
            ],
        ),
        (
            "6. Research Gap",
            [
                "Most solutions stop at classifying an email as spam or ham.",
                "Few methods cluster spam into campaign-level groups.",
                "IOC relationships across domains, IPs, and URLs are underused.",
                "Predictive intelligence from spam behavior is still limited.",
            ],
        ),
        (
            "7. Objectives",
            [
                "Extract IOCs from spam emails (URLs, domains, IPs, sender artifacts).",
                "Cluster related emails into phishing campaigns.",
                "Build a threat intelligence graph to connect related indicators.",
                "Predict likely future attacks from emerging patterns.",
            ],
        ),
        (
            "8. Proposed System Overview",
            [
                "Input: spam and phishing email datasets.",
                "Processing: preprocessing, IOC extraction, and feature engineering.",
                "Analysis: ML clustering and anomaly detection.",
                "Output: threat reports, campaign alerts, and risk scoring.",
            ],
        ),
        (
            "9. Methodology (Part 1)",
            [
                "Email preprocessing: cleaning, normalization, and deduplication.",
                "Feature extraction from text, URLs, sender metadata, and headers.",
                "IOC extraction using regex rules and parser-based validation.",
                "Quality checks to reduce noisy or duplicate indicators.",
            ],
        ),
        (
            "10. Methodology (Part 2)",
            [
                "Machine learning for clustering and anomaly detection.",
                "Campaign detection logic combining content, IOC overlap, and timing.",
                "Graph-based mapping of links between emails, IOCs, and entities.",
                "Advanced add-on: threat graph embeddings/GNN for link prediction.",
                "Advanced add-on: time-aware early warning via change-point detection.",
            ],
        ),
        (
            "12. Tools & Technologies",
            [
                "Python for data processing and automation.",
                "Scikit-learn and NLP libraries for ML modeling.",
                "Graph analytics with NetworkX (and optional Neo4j).",
                "Datasets: public spam/phishing corpora for reproducible experiments.",
            ],
        ),
        (
            "13. Expected Outcome",
            [
                "Earlier detection of emerging phishing campaigns.",
                "Automatically generated IOC intelligence database.",
                "Campaign-level visibility of attacker infrastructure reuse.",
                "Risk-prioritized outputs for SOC investigation workflows.",
            ],
        ),
        (
            "14. Limitations",
            [
                "Performance depends on dataset quality and representativeness.",
                "Public datasets may not fully match enterprise traffic patterns.",
                "System supports analysts; it does not replace full SIEM platforms.",
                "Models require ongoing updates as attacker tactics evolve.",
            ],
        ),
        (
            "15. Conclusion",
            [
                "Threat intelligence from spam data is more valuable than filtering alone.",
                "The proposed pipeline supports early campaign detection and response.",
                "IOC correlation and ML improve visibility of coordinated attacks.",
                "Future scope: real-time SOC integration and continuous learning.",
            ],
        ),
    ]

    for title, bullets in slides[:9]:
        add_bullet_slide(prs, title, bullets)

    add_architecture_slide(prs)

    for title, bullets in slides[9:]:
        add_bullet_slide(prs, title, bullets)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PATH)


if __name__ == "__main__":
    build_presentation()
    print(f"Presentation generated: {OUTPUT_PATH}")
